#!/usr/bin/env python3
"""
将 LLM 生成的内容数据封装成真实二进制文件。
支持格式：PDF, DOCX, DOC, XLSX, XLS, PPTX, PPT, CSV, JSON, XML, TXT

依赖安装：
    pip install reportlab openpyxl python-docx python-pptx xlwt
    sudo apt install libreoffice  # 用于生成 DOC 和 PPT 旧格式

requirements.txt:
    reportlab>=4.0.0
    openpyxl>=3.1.0
    python-docx>=1.0.0
    python-pptx>=0.6.23
    xlwt>=1.3.0
"""

import json
import csv
import xml.etree.ElementTree as ET
import subprocess
from pathlib import Path
from typing import Any, Dict

# 表格类
from openpyxl import Workbook

try:
    import xlwt  # 用于生成 XLS 旧格式

    HAS_XLWT = True
except ImportError:
    HAS_XLWT = False

# 文档类
from reportlab.lib.pagesizes import A4
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
from docx import Document

# 演示文稿类
from pptx import Presentation
from pptx.util import Inches


def generate_pdf(content: Dict, output_path: Path):
    """生成 PDF 文件"""
    doc = SimpleDocTemplate(str(output_path), pagesize=A4)
    elements = []
    styles = getSampleStyleSheet()

    # 标题
    elements.append(Paragraph(content["title"], styles["Title"]))
    elements.append(Spacer(1, 20))

    # 章节和段落
    for section in content.get("sections", []):
        elements.append(Paragraph(section["heading"], styles["Heading2"]))
        for para in section.get("paragraphs", []):
            elements.append(Paragraph(para, styles["Normal"]))
        elements.append(Spacer(1, 10))

    # 表格
    if "table" in content and content["table"]:
        table_data = [content["table"]["headers"]] + content["table"]["rows"]
        table = Table(table_data)
        table.setStyle(
            TableStyle(
                [
                    ("BACKGROUND", (0, 0), (-1, 0), colors.grey),
                    ("TEXTCOLOR", (0, 0), (-1, 0), colors.whitesmoke),
                    ("GRID", (0, 0), (-1, -1), 1, colors.black),
                    ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                    ("ALIGN", (0, 0), (-1, -1), "CENTER"),
                ]
            )
        )
        elements.append(table)

    doc.build(elements)


def generate_docx(content: Dict, output_path: Path):
    """生成 DOCX 文件"""
    doc = Document()
    doc.add_heading(content["title"], 0)

    for section in content.get("sections", []):
        doc.add_heading(section["heading"], level=1)
        for para in section.get("paragraphs", []):
            doc.add_paragraph(para)

    # 表格
    if "table" in content and content["table"]:
        table_data = [content["table"]["headers"]] + content["table"]["rows"]
        table = doc.add_table(len(table_data), len(table_data[0]))
        table.style = "Table Grid"
        for i, row in enumerate(table_data):
            for j, cell in enumerate(row):
                table.cell(i, j).text = str(cell)

    doc.save(str(output_path))


def generate_doc(content: Dict, output_path: Path):
    """生成 DOC 文件（旧格式），通过 LibreOffice 转换"""
    # 先生成 DOCX
    docx_path = output_path.with_suffix(".docx")
    generate_docx(content, docx_path)

    # 使用 LibreOffice 转换为 DOC
    try:
        subprocess.run(
            [
                "libreoffice",
                "--headless",
                "--convert-to",
                "doc",
                str(docx_path),
                "--outdir",
                str(docx_path.parent),
            ],
            check=True,
            capture_output=True,
        )
        # 删除临时 DOCX
        docx_path.unlink()
    except subprocess.CalledProcessError as e:
        raise RuntimeError(f"LibreOffice 转换失败: {e.stderr.decode()}")
    except FileNotFoundError:
        raise RuntimeError("未找到 LibreOffice，请安装: sudo apt install libreoffice")


def generate_xlsx(content: Dict, output_path: Path):
    """生成 XLSX 文件"""
    wb = Workbook()

    for idx, sheet_data in enumerate(content.get("sheets", [])):
        if idx == 0:
            ws = wb.active
            ws.title = sheet_data.get("name", "Sheet1")
        else:
            ws = wb.create_sheet(sheet_data.get("name", f"Sheet{idx+1}"))

        ws.append(sheet_data["headers"])
        for row in sheet_data.get("rows", []):
            ws.append(row)

    wb.save(str(output_path))


def generate_xls(content: Dict, output_path: Path):
    """生成 XLS 文件（旧格式）"""
    if not HAS_XLWT:
        raise ImportError("需要安装 xlwt: pip install xlwt")

    wb = xlwt.Workbook()

    for idx, sheet_data in enumerate(content.get("sheets", [])):
        ws = wb.add_sheet(sheet_data.get("name", f"Sheet{idx+1}"))

        # 表头样式
        header_style = xlwt.XFStyle()
        pattern = xlwt.Pattern()
        pattern.pattern = xlwt.Pattern.SOLID_PATTERN
        pattern.pattern_fore_colour = xlwt.Style.colour_map["grey25"]
        header_style.pattern = pattern

        # 写入表头
        for j, header in enumerate(sheet_data["headers"]):
            ws.write(0, j, header, header_style)

        # 写入数据
        for i, row in enumerate(sheet_data.get("rows", []), start=1):
            for j, cell in enumerate(row):
                ws.write(i, j, cell)

    wb.save(str(output_path))


def generate_pptx(content: Dict, output_path: Path):
    """生成 PPTX 文件"""
    prs = Presentation()

    for slide_data in content.get("slides", []):
        layout = slide_data.get("layout", "content")

        if layout == "title":
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = slide_data.get("title", "")
            if "subtitle" in slide_data:
                slide.placeholders[1].text = slide_data["subtitle"]

        elif layout == "content":
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = slide_data.get("title", "")
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()
            for i, bullet in enumerate(slide_data.get("bullets", [])):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = bullet
                p.level = 0

        elif layout == "table":
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = slide_data.get("title", "")
            table_data = [slide_data["table"]["headers"]] + slide_data["table"]["rows"]
            rows = len(table_data)
            cols = len(table_data[0]) if table_data else 0
            table = slide.shapes.add_table(
                rows, cols, Inches(1), Inches(2), Inches(8), Inches(4)
            ).table
            for i, row in enumerate(table_data):
                for j, cell in enumerate(row):
                    table.cell(i, j).text = str(cell)

    prs.save(str(output_path))


def generate_ppt(content: Dict, output_path: Path):
    """生成 PPT 文件（旧格式），通过 LibreOffice 转换"""
    # 先生成 PPTX
    pptx_path = output_path.with_suffix(".pptx")
    generate_pptx(content, pptx_path)

    # 使用 LibreOffice 转换为 PPT
    try:
        subprocess.run(
            [
                "libreoffice",
                "--headless",
                "--convert-to",
                "ppt",
                str(pptx_path),
                "--outdir",
                str(pptx_path.parent),
            ],
            check=True,
            capture_output=True,
        )
        # 删除临时 PPTX
        pptx_path.unlink()
    except subprocess.CalledProcessError as e:
        raise RuntimeError(f"LibreOffice 转换失败: {e.stderr.decode()}")
    except FileNotFoundError:
        raise RuntimeError("未找到 LibreOffice，请安装: sudo apt install libreoffice")


def generate_csv(content: Dict, output_path: Path):
    """生成 CSV 文件"""
    with open(str(output_path), "w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        writer.writerow(content["headers"])
        writer.writerows(content["rows"])


def generate_json(content: Dict, output_path: Path):
    """生成 JSON 文件"""
    with open(str(output_path), "w", encoding="utf-8") as f:
        json.dump(content["data"], f, ensure_ascii=False, indent=2)


def generate_xml(content: Dict, output_path: Path):
    """生成 XML 文件"""

    def build_element(parent: ET.Element, elem_data: Dict):
        elem = ET.SubElement(parent, elem_data["tag"])
        for k, v in elem_data.get("attrs", {}).items():
            elem.set(k, v)
        if "text" in elem_data:
            elem.text = elem_data["text"]
        for child in elem_data.get("children", []):
            build_element(elem, child)

    root = ET.Element(content["root"])
    for elem_data in content.get("elements", []):
        build_element(root, elem_data)

    tree = ET.ElementTree(root)
    ET.indent(tree, space="  ", level=0)
    tree.write(str(output_path), encoding="utf-8", xml_declaration=True)


def generate_txt(content: Dict, output_path: Path):
    """生成 TXT 文件"""
    with open(str(output_path), "w", encoding="utf-8") as f:
        f.write("\n".join(content["lines"]))


# 格式到生成函数的映射
format_handlers = {
    "pdf": generate_pdf,
    "docx": generate_docx,
    "doc": generate_doc,
    "xlsx": generate_xlsx,
    "xls": generate_xls,
    "pptx": generate_pptx,
    "ppt": generate_ppt,
    "csv": generate_csv,
    "json": generate_json,
    "xml": generate_xml,
    "txt": generate_txt,
}


def generate_file(content_json: Dict, output_dir: str = "test_files") -> Path:
    """
    根据 LLM 生成的内容数据生成文件。

    Args:
        content_json: LLM 返回的 JSON 数据，包含 filename 和 content
        output_dir: 输出目录

    Returns:
        生成的文件路径
    """
    output_path = Path(output_dir) / content_json["filename"]
    output_path.parent.mkdir(parents=True, exist_ok=True)

    # 获取文件格式
    file_format = output_path.suffix.lstrip(".").lower()

    # 获取对应的生成函数
    handler = format_handlers.get(file_format)
    if not handler:
        raise ValueError(f"不支持的文件格式：{file_format}")

    # 生成文件
    handler(content_json["content"], output_path)
    return output_path


def generate_files_from_json_list(
    files_json: list, output_dir: str = "test_files"
) -> list:
    """
    批量生成文件。

    Args:
        files_json: LLM 返回的文件列表
        output_dir: 输出目录

    Returns:
        生成的文件路径列表
    """
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    generated = []
    for file_info in files_json:
        try:
            path = generate_file(file_info, str(output_dir))
            generated.append(path)
            print(f"✓ 生成：{path}")
        except Exception as e:
            print(f"✗ 失败 {file_info.get('filename', 'unknown')}: {e}")

    return generated


def main():
    """从 llm_content.json 读取并生成所有文件"""
    import sys

    # 默认从 llm_content.json 读取
    input_file = sys.argv[1] if len(sys.argv) > 1 else "llm_content.json"
    output_dir = sys.argv[2] if len(sys.argv) > 2 else "test_files"

    with open(input_file, "r", encoding="utf-8") as f:
        data = json.load(f)

    # 支持两种格式：单个文件或文件列表
    if "files" in data:
        files = data["files"]
    elif isinstance(data, list):
        files = data
    else:
        files = [data]

    generated = generate_files_from_json_list(files, output_dir)
    print(f"\n共生成 {len(generated)} 个文件")


if __name__ == "__main__":
    main()
